import fitz
from pptx import Presentation

def extraer_texto_pdf(path):
    doc = fitz.open(path)
    return "\n".join([page.get_text() for page in doc])

def extraer_texto_ppt(path):
    prs = Presentation(path)
    texto = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texto.append(shape.text)
    return "\n".join(texto)
